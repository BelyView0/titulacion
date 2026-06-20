import io
from django.http import HttpResponse
from django.views import View
from django.db.models import Count, Q
from expediente.mixins import JefeProyectoRequeridoMixin
from expediente.models import Expediente, EstadoExpediente
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import PieChart, BarChart, Reference
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def get_estadisticas_data(user):
    departamento = user.departamento
    if departamento:
        qs_base = Expediente.objects.filter(alumno__carrera__departamento=departamento)
    else:
        qs_base = Expediente.objects.filter(alumno__carrera=user.carrera)

    qs = qs_base.select_related('alumno', 'modalidad')

    total_expedientes = qs.count()
    qs_concluidos = qs.filter(estado=EstadoExpediente.CONCLUIDO)
    total_concluidos = qs_concluidos.count()
    expedientes_activos = qs.exclude(
        estado__in=[EstadoExpediente.CONCLUIDO, EstadoExpediente.CANCELADO, EstadoExpediente.BORRADOR]
    ).count()
    expedientes_cancelados = qs.filter(estado=EstadoExpediente.CANCELADO).count()

    porcentaje_titulados = round(
        (total_concluidos / total_expedientes * 100) if total_expedientes > 0 else 0, 1
    )

    iniciados_hombres = qs.filter(alumno__genero='M').count()
    iniciados_mujeres = qs.filter(alumno__genero='F').count()
    iniciados_sin_dato = total_expedientes - iniciados_hombres - iniciados_mujeres

    titulados_hombres = qs_concluidos.filter(alumno__genero='M').count()
    titulados_mujeres = qs_concluidos.filter(alumno__genero='F').count()
    titulados_sin_dato = total_concluidos - titulados_hombres - titulados_mujeres

    por_generacion = list(
        qs.filter(alumno__generacion__isnull=False)
        .values('alumno__generacion')
        .annotate(
            total=Count('id'),
            concluidos=Count('id', filter=Q(estado=EstadoExpediente.CONCLUIDO))
        )
        .order_by('-alumno__generacion')
    )

    por_modalidad = list(
        qs.filter(modalidad__isnull=False)
        .values('modalidad__nombre')
        .annotate(
            total=Count('id'),
            concluidos=Count('id', filter=Q(estado=EstadoExpediente.CONCLUIDO))
        )
        .order_by('-total')
    )

    por_carrera = list(
        qs.values('alumno__carrera__nombre')
        .annotate(
            total=Count('id'),
            concluidos=Count('id', filter=Q(estado=EstadoExpediente.CONCLUIDO)),
            hombres=Count('id', filter=Q(alumno__genero='M')),
            mujeres=Count('id', filter=Q(alumno__genero='F')),
        )
        .order_by('-total')
    )

    estado_display = dict(EstadoExpediente.choices)
    expedientes_por_estado = [
        {
            'estado': estado_display.get(item['estado'], item['estado']),
            'total': item['total'],
            'clave': item['estado'],
        }
        for item in (
            qs.values('estado')
            .annotate(total=Count('id'))
            .order_by('-total')
        )
    ]

    return {
        'departamento': departamento.nombre if departamento else (user.carrera.nombre if user.carrera else 'Global'),
        'total_expedientes': total_expedientes,
        'total_concluidos': total_concluidos,
        'expedientes_activos': expedientes_activos,
        'expedientes_cancelados': expedientes_cancelados,
        'porcentaje_titulados': porcentaje_titulados,
        'iniciados_hombres': iniciados_hombres,
        'iniciados_mujeres': iniciados_mujeres,
        'iniciados_sin_dato': iniciados_sin_dato,
        'titulados_hombres': titulados_hombres,
        'titulados_mujeres': titulados_mujeres,
        'titulados_sin_dato': titulados_sin_dato,
        'por_generacion': por_generacion,
        'por_modalidad': por_modalidad,
        'por_carrera': por_carrera,
        'expedientes_por_estado': expedientes_por_estado,
    }


class ExportarEstadisticasExcelView(JefeProyectoRequeridoMixin, View):
    def get(self, request, *args, **kwargs):
        stats = get_estadisticas_data(request.user)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Resumen General"

        # Estilos
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="1B396A", end_color="1B396A", fill_type="solid")
        title_font = Font(size=14, bold=True)

        ws['A1'] = f"Estadísticas de Titulación - {stats['departamento']}"
        ws['A1'].font = title_font
        ws.merge_cells('A1:C1')
        ws.append([])
        
        # Resumen Global
        ws.append(["Métrica", "Cantidad"])
        for cell in ws[3]:
            cell.font = header_font; cell.fill = header_fill

        ws.append(["Total de Expedientes Iniciados", stats['total_expedientes']])
        ws.append(["Total Titulados (Concluidos)", stats['total_concluidos']])
        ws.append(["Expedientes en Proceso (Activos)", stats['expedientes_activos']])
        ws.append(["Expedientes Cancelados", stats['expedientes_cancelados']])
        ws.append(["Eficiencia de Titulación", f"{stats['porcentaje_titulados']}%"])
        ws.append([])

        # Estatus (Tabla para Gráfica de Pastel)
        row_start_estados = ws.max_row + 1
        ws.append(["Estado del Expediente", "Total"])
        for cell in ws[row_start_estados]:
            cell.font = header_font; cell.fill = header_fill
        
        for est in stats['expedientes_por_estado']:
            ws.append([est['estado'], est['total']])
        row_end_estados = ws.max_row
        ws.append([])

        # Gráfica de Pastel - Estados
        if len(stats['expedientes_por_estado']) > 0:
            pie = PieChart()
            labels = Reference(ws, min_col=1, min_row=row_start_estados+1, max_row=row_end_estados)
            data = Reference(ws, min_col=2, min_row=row_start_estados, max_row=row_end_estados)
            pie.add_data(data, titles_from_data=True)
            pie.set_categories(labels)
            pie.title = "Distribución por Estatus"
            ws.add_chart(pie, "D3")

        # Género
        ws.append(["Género", "Iniciaron Proceso", "Concluyeron Titulación"])
        for cell in ws[ws.max_row]:
            cell.font = header_font; cell.fill = header_fill
        ws.append(["Hombres", stats['iniciados_hombres'], stats['titulados_hombres']])
        ws.append(["Mujeres", stats['iniciados_mujeres'], stats['titulados_mujeres']])
        ws.append(["Sin especificar", stats['iniciados_sin_dato'], stats['titulados_sin_dato']])
        ws.append([])

        # Generación
        ws.append(["Generación", "Total Iniciaron", "Total Titulados"])
        for cell in ws[ws.max_row]:
            cell.font = header_font; cell.fill = header_fill
        for gen in stats['por_generacion']:
            ws.append([gen['alumno__generacion'], gen['total'], gen['concluidos']])

        # Hoja 2: Modalidades y Carreras
        ws2 = wb.create_sheet(title="Modalidades y Carreras")
        
        row_start_mod = ws2.max_row
        ws2.append(["Modalidad", "Total Iniciaron", "Total Titulados"])
        for cell in ws2[1]:
            cell.font = header_font; cell.fill = header_fill
        for mod in stats['por_modalidad']:
            ws2.append([mod['modalidad__nombre'], mod['total'], mod['concluidos']])
        row_end_mod = ws2.max_row

        # Gráfica de Barras - Modalidades
        if len(stats['por_modalidad']) > 0:
            bar = BarChart()
            bar.type = "col"
            bar.style = 10
            bar.title = "Titulados por Modalidad"
            labels_mod = Reference(ws2, min_col=1, min_row=row_start_mod+1, max_row=row_end_mod)
            data_mod = Reference(ws2, min_col=3, min_row=row_start_mod, max_row=row_end_mod)
            bar.add_data(data_mod, titles_from_data=True)
            bar.set_categories(labels_mod)
            ws2.add_chart(bar, "E2")

        ws2.append([])
        ws2.append(["Carrera", "Total Iniciaron", "Total Titulados", "Hombres", "Mujeres"])
        for cell in ws2[ws2.max_row]:
            cell.font = header_font; cell.fill = header_fill
        for car in stats['por_carrera']:
            ws2.append([car['alumno__carrera__nombre'], car['total'], car['concluidos'], car['hombres'], car['mujeres']])

        # Ajustar ancho columnas
        from openpyxl.utils import get_column_letter
        for sheet in wb.worksheets:
            for col in sheet.columns:
                max_length = 0
                col_letter = get_column_letter(col[0].column)
                for cell in col:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                sheet.column_dimensions[col_letter].width = min(max_length + 2, 60)

        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        response['Content-Disposition'] = f'attachment; filename="Estadisticas_Titulacion.xlsx"'
        wb.save(response)
        return response


class ExportarEstadisticasPPTXView(JefeProyectoRequeridoMixin, View):
    def get(self, request, *args, **kwargs):
        stats = get_estadisticas_data(request.user)

        prs = Presentation()
        
        # Helper colors
        BLUE_DARK = RGBColor(27, 57, 106)
        BLUE_LIGHT = RGBColor(0, 114, 198)
        GRAY = RGBColor(108, 117, 125)

        # Portada
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Reporte Ejecutivo de Titulación"
        title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = f"Departamento: {stats['departamento']}\nSistema HEDS"
        subtitle.text_frame.paragraphs[0].font.color.rgb = GRAY

        # Diapositiva 2: Resumen
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide2.shapes.title
        title_shape.text = "Resumen General"
        title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
        
        tf = slide2.shapes.placeholders[1].text_frame
        tf.text = f"Total de Expedientes: {stats['total_expedientes']}"
        tf.add_paragraph().text = f"Total de Titulados: {stats['total_concluidos']} ({stats['porcentaje_titulados']}%)"
        tf.add_paragraph().text = f"Expedientes en Proceso: {stats['expedientes_activos']}"
        tf.add_paragraph().text = f"Expedientes Cancelados: {stats['expedientes_cancelados']}"

        # Diapositiva 3: Gráfica de Estados (Pastel)
        blank_slide_layout = prs.slide_layouts[5] # Título solamente
        slide_estados = prs.slides.add_slide(blank_slide_layout)
        slide_estados.shapes.title.text = "Distribución por Estatus del Expediente"
        slide_estados.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK

        if stats['expedientes_por_estado']:
            chart_data = CategoryChartData()
            chart_data.categories = [item['estado'] for item in stats['expedientes_por_estado']]
            chart_data.add_series('Total', (item['total'] for item in stats['expedientes_por_estado']))

            x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
            chart = slide_estados.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.plots[0].has_data_labels = True

        # Diapositiva 4: Demografía (Barras apiladas)
        slide_demo = prs.slides.add_slide(blank_slide_layout)
        slide_demo.shapes.title.text = "Demografía: Hombres vs Mujeres"
        slide_demo.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
        
        chart_data_demo = CategoryChartData()
        chart_data_demo.categories = ['Iniciaron Proceso', 'Titulados (Concluidos)']
        chart_data_demo.add_series('Hombres', (stats['iniciados_hombres'], stats['titulados_hombres']))
        chart_data_demo.add_series('Mujeres', (stats['iniciados_mujeres'], stats['titulados_mujeres']))
        if stats['iniciados_sin_dato'] > 0 or stats['titulados_sin_dato'] > 0:
            chart_data_demo.add_series('No especificado', (stats['iniciados_sin_dato'], stats['titulados_sin_dato']))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
        chart_demo = slide_demo.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_demo
        ).chart
        chart_demo.has_legend = True
        chart_demo.legend.position = XL_LEGEND_POSITION.BOTTOM

        # Diapositiva 5: Gráfica de Modalidades (Barras)
        slide_mod = prs.slides.add_slide(blank_slide_layout)
        slide_mod.shapes.title.text = "Top Modalidades de Titulación"
        slide_mod.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK

        if stats['por_modalidad']:
            chart_data_mod = CategoryChartData()
            # Tomar top 5
            top_mod = stats['por_modalidad'][:5]
            chart_data_mod.categories = [item['modalidad__nombre'][:25] + '...' if len(item['modalidad__nombre']) > 25 else item['modalidad__nombre'] for item in top_mod]
            chart_data_mod.add_series('Concluidos', (item['concluidos'] for item in top_mod))
            chart_data_mod.add_series('Iniciados', (item['total'] for item in top_mod))

            x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(4.5)
            chart_mod = slide_mod.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data_mod
            ).chart
            chart_mod.has_legend = True
            chart_mod.legend.position = XL_LEGEND_POSITION.BOTTOM
            
        # Diapositiva 6: Gráfica por Carrera
        if stats['por_carrera']:
            slide_car = prs.slides.add_slide(blank_slide_layout)
            slide_car.shapes.title.text = "Expedientes por Carrera"
            slide_car.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
            
            chart_data_car = CategoryChartData()
            chart_data_car.categories = [item['alumno__carrera__nombre'][:25] + '...' if len(item['alumno__carrera__nombre']) > 25 else item['alumno__carrera__nombre'] for item in stats['por_carrera']]
            chart_data_car.add_series('Iniciados', (item['total'] for item in stats['por_carrera']))
            chart_data_car.add_series('Concluidos', (item['concluidos'] for item in stats['por_carrera']))
            
            x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(4.5)
            chart_car = slide_car.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_car
            ).chart
            chart_car.has_legend = True
            chart_car.legend.position = XL_LEGEND_POSITION.BOTTOM

        # Diapositiva 7: Gráfica por Generación
        if stats['por_generacion']:
            slide_gen = prs.slides.add_slide(blank_slide_layout)
            slide_gen.shapes.title.text = "Expedientes por Generación"
            slide_gen.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
            
            chart_data_gen = CategoryChartData()
            top_gen = stats['por_generacion'][:10] # Mostrar las últimas 10 generaciones
            top_gen.reverse() # Cronológico de izquierda a derecha
            chart_data_gen.categories = [str(item['alumno__generacion']) for item in top_gen]
            chart_data_gen.add_series('Iniciados', (item['total'] for item in top_gen))
            chart_data_gen.add_series('Concluidos', (item['concluidos'] for item in top_gen))
            
            x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(4.5)
            chart_gen = slide_gen.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_gen
            ).chart
            chart_gen.has_legend = True
            chart_gen.legend.position = XL_LEGEND_POSITION.BOTTOM

        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = f'attachment; filename="Presentacion_Estadisticas.pptx"'
        prs.save(response)
        return response
